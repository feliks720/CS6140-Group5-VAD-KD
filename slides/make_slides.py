"""Generate presentation slides for VAD Knowledge Distillation project."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color palette ──
BG_DARK   = RGBColor(0x1B, 0x1B, 0x2F)   # deep navy
BG_CARD   = RGBColor(0x26, 0x26, 0x40)   # card bg
ACCENT    = RGBColor(0x4E, 0xC9, 0xB0)   # teal green
ACCENT2   = RGBColor(0x56, 0x9C, 0xD6)   # blue
ACCENT3   = RGBColor(0xDC, 0xDC, 0xAA)   # gold
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0xAA, 0xAA, 0xBB)
LIGHT     = RGBColor(0xDD, 0xDD, 0xEE)
RED_SOFT  = RGBColor(0xF4, 0x72, 0x6C)
GREEN_OK  = RGBColor(0x6A, 0xCE, 0x75)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    adj = shape.adjustments
    if len(adj) > 0:
        adj[0] = 0.04
    return shape

def add_bullet_box(slide, left, top, width, height, items, font_size=16,
                   color=LIGHT, bullet_color=ACCENT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(8)
        # Bullet character
        run_b = p.add_run()
        run_b.text = "  >  " if not item.startswith("!") else "  "
        run_b.font.size = Pt(font_size)
        run_b.font.color.rgb = bullet_color
        run_b.font.bold = True
        run_b.font.name = "Consolas"
        # Text
        clean = item.lstrip("!")
        run_t = p.add_run()
        run_t.text = clean
        run_t.font.size = Pt(font_size)
        run_t.font.color.rgb = color
        run_t.font.name = "Segoe UI"
    return txBox

def add_metric_card(slide, left, top, width, value, label, val_color=ACCENT):
    add_rounded_rect(slide, left, top, width, 1.15, BG_CARD)
    add_text_box(slide, left, top + 0.08, width, 0.6, value,
                 font_size=28, color=val_color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + 0.62, width, 0.45, label,
                 font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 1 — Title                                           ║
# ╚══════════════════════════════════════════════════════════════╝
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BG_DARK)

# Decorative line
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(1.5), Inches(2.4), Inches(3.5), Inches(0.06))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

add_text_box(slide, 1.5, 1.0, 10, 1.4,
    "Knowledge Distillation for\nVoice Activity Detection",
    font_size=42, color=WHITE, bold=True)
add_text_box(slide, 1.5, 2.7, 10, 0.6,
    "Compressing CRDNN Teacher into Lightweight Student Models",
    font_size=22, color=ACCENT)
add_text_box(slide, 1.5, 3.6, 10, 0.5,
    "CS6140 Machine Learning  |  Northeastern University  |  Group 5",
    font_size=16, color=GRAY)

# Right side - key numbers preview
add_rounded_rect(slide, 8.5, 1.2, 3.8, 5.0, BG_CARD)
add_text_box(slide, 8.5, 1.4, 3.8, 0.5, "Highlights",
             font_size=18, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

highlights = [
    "3 student architectures evaluated",
    "Best F1: 0.8122 (TinyTransformer)",
    "7.4x compression (TinyCNN: 57 KB)",
    "Temperature sweep: T = {1, 2, 4, 8}",
    "Dataset: LibriParty (350 sessions)",
]
add_bullet_box(slide, 8.8, 2.0, 3.3, 3.8, highlights, font_size=14, color=LIGHT)


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 2 — Problem & Motivation                            ║
# ╚══════════════════════════════════════════════════════════════╝
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, 0.8, 0.35, 6, 0.6, "Problem & Motivation", font_size=32, color=WHITE, bold=True)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.8), Inches(0.95), Inches(2.5), Inches(0.05))
shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()

# Left column - Problem
add_rounded_rect(slide, 0.8, 1.3, 5.6, 5.6, BG_CARD)
add_text_box(slide, 1.1, 1.45, 5, 0.5, "The Challenge", font_size=20, color=ACCENT, bold=True)
items_left = [
    "Neural VAD models (CRDNN) achieve high accuracy but are too large for edge devices",
    "Energy-based VAD is lightweight but fails in noisy, multi-speaker audio",
    "Need: compact models that retain neural-level quality",
    "!",
    "!Knowledge Distillation transfers dark knowledge",
    "!from a large teacher to a small student via",
    "!softened probability distributions",
]
add_bullet_box(slide, 1.1, 2.1, 5.0, 4.5, items_left, font_size=15)

# Right column - Evidence
add_rounded_rect(slide, 6.8, 1.3, 5.6, 2.5, BG_CARD)
add_text_box(slide, 7.1, 1.45, 5, 0.5, "Energy VAD vs Neural VAD", font_size=20, color=ACCENT2, bold=True)

add_metric_card(slide, 7.2, 2.1, 2.3, "7.4%", "Energy VAD speech ratio", val_color=RED_SOFT)
add_metric_card(slide, 9.8, 2.1, 2.3, "45.6%", "Teacher speech ratio", val_color=GREEN_OK)

# KD formula box
add_rounded_rect(slide, 6.8, 4.1, 5.6, 2.8, BG_CARD)
add_text_box(slide, 7.1, 4.25, 5, 0.5, "KD Loss Function", font_size=20, color=ACCENT2, bold=True)
add_text_box(slide, 7.1, 4.85, 5.0, 0.7,
    "L = a * KL(teacher_soft || student_soft) * T^2\n    + (1-a) * BCE(student, hard_label)",
    font_size=16, color=ACCENT3, font_name="Consolas")
kd_items = [
    "a=0.7: KD loss dominates (70% soft, 30% hard)",
    "T: temperature controlling softness of distributions",
    "T^2 scaling preserves gradient magnitudes",
]
add_bullet_box(slide, 7.1, 5.7, 5.0, 1.5, kd_items, font_size=13, color=GRAY)


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 3 — Student Architectures                           ║
# ╚══════════════════════════════════════════════════════════════╝
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, 0.8, 0.35, 8, 0.6, "Student Architectures", font_size=32, color=WHITE, bold=True)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.8), Inches(0.95), Inches(2.5), Inches(0.05))
shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()

# Three architecture cards
models = [
    ("TinyCNN", "14,913", "57 KB",
     ["3x (Conv1D + BN + ReLU + Dropout)", "Kernel sizes: 5, 5, 3", "Captures local spectral patterns", "No recurrence - fully parallel"]),
    ("MLP", "81,921", "313 KB",
     ["Context window: 5 frames each side", "3-layer MLP per frame", "LayerNorm + ReLU + Dropout", "Frame-independent processing"]),
    ("TinyTransformer", "389,633", "1.49 MB",
     ["2-layer self-attention (2 heads)", "d_model=64, d_ff=128", "Sinusoidal positional encoding", "Global temporal modeling"]),
]

for i, (name, params, size, desc) in enumerate(models):
    x = 0.8 + i * 4.1
    add_rounded_rect(slide, x, 1.3, 3.8, 5.6, BG_CARD)
    add_text_box(slide, x, 1.45, 3.8, 0.5, name,
                 font_size=22, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

    add_metric_card(slide, x + 0.2, 2.1, 1.5, params, "Parameters",
                    val_color=[ACCENT, ACCENT2, ACCENT3][i])
    add_metric_card(slide, x + 2.0, 2.1, 1.5, size, "Model Size",
                    val_color=[ACCENT, ACCENT2, ACCENT3][i])

    add_bullet_box(slide, x + 0.3, 3.6, 3.2, 3.0, desc, font_size=14, color=LIGHT)

# Teacher reference
add_text_box(slide, 0.8, 7.0, 12, 0.4,
    "Teacher (CRDNN):  109,744 params  |  0.43 MB  |  SpeechBrain pretrained on LibriParty",
    font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 4 — Main Results Table                              ║
# ╚══════════════════════════════════════════════════════════════╝
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, 0.8, 0.35, 10, 0.6, "Results: Student Architecture Comparison",
             font_size=32, color=WHITE, bold=True)
add_text_box(slide, 0.8, 0.85, 10, 0.4, "T = 4,  alpha = 0.7,  30 epochs,  evaluated on dev set",
             font_size=16, color=GRAY)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.8), Inches(1.2), Inches(2.5), Inches(0.05))
shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()

# Table
headers = ["Model", "Params", "Size", "F1", "Precision", "Recall", "Accuracy", "DER"]
rows = [
    ["Energy VAD",       "0",       "N/A",     "--",     "--",     "--",     "--",     "--"],
    ["TinyCNN",          "14,913",  "57 KB",   "0.7707", "0.7774", "0.7640", "0.8270", "0.3703"],
    ["MLP",              "81,921",  "313 KB",  "0.7841", "0.7651", "0.8040", "0.8315", "0.3476"],
    ["TinyTransformer",  "389,633", "1.49 MB", "0.8122", "0.7542", "0.8800", "0.8452", "0.2962"],
]

table_shape = slide.shapes.add_table(len(rows) + 1, len(headers),
    Inches(0.8), Inches(1.5), Inches(11.7), Inches(2.8))
table = table_shape.table

# Style header
for j, h in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = h
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.bold = True
        paragraph.font.color.rgb = BG_DARK
        paragraph.font.name = "Segoe UI"
        paragraph.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = ACCENT

# Fill rows
best_row = 3  # TinyTransformer
for i, row in enumerate(rows):
    for j, val in enumerate(row):
        cell = table.cell(i + 1, j)
        cell.text = val
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(13)
            paragraph.font.name = "Segoe UI"
            paragraph.alignment = PP_ALIGN.CENTER
            if i == best_row and j >= 3:
                paragraph.font.color.rgb = ACCENT
                paragraph.font.bold = True
            elif i == 0:
                paragraph.font.color.rgb = RED_SOFT
            else:
                paragraph.font.color.rgb = LIGHT
        cell.fill.solid()
        cell.fill.fore_color.rgb = BG_CARD if i % 2 == 0 else RGBColor(0x2E, 0x2E, 0x4A)

# Observations
add_rounded_rect(slide, 0.8, 4.6, 5.6, 2.6, BG_CARD)
add_text_box(slide, 1.1, 4.7, 5, 0.4, "Key Findings", font_size=18, color=ACCENT, bold=True)
findings = [
    "TinyTransformer: best quality (F1=0.81), self-attention captures global context",
    "TinyCNN: best efficiency -- 7.4x compression, only 57 KB",
    "MLP: balanced trade-off (F1=0.78, 313 KB)",
    "All students >> energy VAD baseline",
]
add_bullet_box(slide, 1.1, 5.2, 5.0, 2.2, findings, font_size=13, color=LIGHT)

# Compression visual
add_rounded_rect(slide, 6.8, 4.6, 5.6, 2.6, BG_CARD)
add_text_box(slide, 7.1, 4.7, 5, 0.4, "Compression vs Teacher (109K params)",
             font_size=18, color=ACCENT2, bold=True)
comp_items = [
    "TinyCNN:          7.4x smaller   (14.9K params)",
    "MLP:              1.3x smaller   (81.9K params)",
    "TinyTransformer:  3.5x LARGER    (389K params)",
    "!",
    "!TinyCNN is the practical edge-deployment choice",
]
add_bullet_box(slide, 7.1, 5.2, 5.0, 2.2, comp_items, font_size=13, color=LIGHT, bullet_color=ACCENT2)


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 5 — Temperature Sweep                               ║
# ╚══════════════════════════════════════════════════════════════╝
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, 0.8, 0.35, 10, 0.6, "Temperature Sweep: Effect of T on TinyCNN",
             font_size=32, color=WHITE, bold=True)
add_text_box(slide, 0.8, 0.85, 10, 0.4, "Student = TinyCNN,  alpha = 0.7,  30 epochs",
             font_size=16, color=GRAY)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.8), Inches(1.2), Inches(2.5), Inches(0.05))
shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()

# Table
t_headers = ["Temperature", "F1", "Precision", "Recall", "Accuracy", "DER", "Best Epoch"]
t_rows = [
    ["T = 1", "0.7759", "0.8104", "0.7443", "0.8365", "0.3626", "17"],
    ["T = 2", "0.7720", "0.7634", "0.7808", "0.8245", "0.3678", "23"],
    ["T = 4", "0.7707", "0.7774", "0.7640", "0.8270", "0.3703", "24"],
    ["T = 8", "0.7667", "0.7595", "0.7741", "0.8208", "0.3764", "22"],
]

table_shape = slide.shapes.add_table(len(t_rows) + 1, len(t_headers),
    Inches(0.8), Inches(1.5), Inches(11.7), Inches(2.6))
table = table_shape.table

for j, h in enumerate(t_headers):
    cell = table.cell(0, j)
    cell.text = h
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.bold = True
        paragraph.font.color.rgb = BG_DARK
        paragraph.font.name = "Segoe UI"
        paragraph.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = ACCENT3

for i, row in enumerate(t_rows):
    for j, val in enumerate(row):
        cell = table.cell(i + 1, j)
        cell.text = val
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.name = "Segoe UI"
            paragraph.alignment = PP_ALIGN.CENTER
            if i == 0 and j >= 1:
                paragraph.font.color.rgb = ACCENT3
                paragraph.font.bold = True
            else:
                paragraph.font.color.rgb = LIGHT
        cell.fill.solid()
        cell.fill.fore_color.rgb = BG_CARD if i % 2 == 0 else RGBColor(0x2E, 0x2E, 0x4A)

# Bar chart simulation using shapes
add_rounded_rect(slide, 0.8, 4.4, 5.6, 2.8, BG_CARD)
add_text_box(slide, 1.1, 4.5, 5, 0.4, "F1 Score by Temperature", font_size=18, color=ACCENT3, bold=True)

bar_data = [("T=1", 0.7759), ("T=2", 0.7720), ("T=4", 0.7707), ("T=8", 0.7667)]
bar_min, bar_max = 0.76, 0.78
bar_base_y = 6.8
bar_max_h = 1.5

for idx, (label, val) in enumerate(bar_data):
    x = 1.5 + idx * 1.1
    h = (val - bar_min) / (bar_max - bar_min) * bar_max_h
    h = max(h, 0.15)
    y = bar_base_y - h
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(0.7), Inches(h))
    bar.fill.solid()
    c = [ACCENT3, ACCENT2, ACCENT, GRAY][idx]
    bar.fill.fore_color.rgb = c
    bar.line.fill.background()
    add_text_box(slide, x - 0.1, bar_base_y + 0.02, 0.9, 0.3, label,
                 font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x - 0.15, y - 0.3, 1.0, 0.3, f"{val:.4f}",
                 font_size=11, color=c, bold=True, alignment=PP_ALIGN.CENTER)

# Analysis
add_rounded_rect(slide, 6.8, 4.4, 5.6, 2.8, BG_CARD)
add_text_box(slide, 7.1, 4.5, 5, 0.4, "Analysis", font_size=18, color=ACCENT2, bold=True)
analysis = [
    "Lower T yields slightly better F1 for TinyCNN",
    "T=1: highest precision (0.81), fastest convergence (epoch 17)",
    "Higher T shifts trade-off toward recall at cost of precision",
    "Small F1 range (0.77) suggests model capacity is the bottleneck",
    "Binary VAD may benefit less from soft labels than multi-class tasks",
]
add_bullet_box(slide, 7.1, 5.0, 5.0, 2.5, analysis, font_size=13, color=LIGHT, bullet_color=ACCENT2)


# ╔══════════════════════════════════════════════════════════════╗
# ║  SLIDE 6 — Conclusions                                     ║
# ╚══════════════════════════════════════════════════════════════╝
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_text_box(slide, 0.8, 0.35, 10, 0.6, "Conclusions", font_size=32, color=WHITE, bold=True)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.8), Inches(0.95), Inches(2.5), Inches(0.05))
shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()

# Main conclusions
add_rounded_rect(slide, 0.8, 1.3, 7.6, 3.0, BG_CARD)
add_text_box(slide, 1.1, 1.4, 7, 0.4, "Key Takeaways", font_size=20, color=ACCENT, bold=True)
conclusions = [
    "KD successfully transfers neural VAD quality to lightweight students, far surpassing energy-based VAD",
    "TinyCNN achieves 7.4x compression at F1=0.77 -- suitable for real-time edge deployment (57 KB)",
    "TinyTransformer reaches the best F1=0.81, but at 3.5x MORE parameters than the teacher",
    "Temperature has minor impact on binary VAD -- T=1 slightly best, suggesting hard teacher decisions are already informative",
    "The precision-recall trade-off is tunable: TinyCNN favors precision, TinyTransformer favors recall",
]
add_bullet_box(slide, 1.1, 1.9, 7.0, 2.5, conclusions, font_size=14, color=LIGHT)

# Best model cards
add_text_box(slide, 8.8, 1.3, 4, 0.4, "Best Results", font_size=18, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
add_rounded_rect(slide, 8.8, 1.8, 3.8, 2.5, BG_CARD)
cards = [
    ("Best Quality", "TinyTransformer", "F1 = 0.8122", ACCENT3),
    ("Best Efficiency", "TinyCNN", "14,913 params", ACCENT),
    ("Best Balance", "MLP", "F1=0.78 / 313 KB", ACCENT2),
]
for i, (title, model, val, c) in enumerate(cards):
    y = 1.95 + i * 0.75
    add_text_box(slide, 9.0, y, 1.6, 0.35, title, font_size=11, color=GRAY)
    add_text_box(slide, 9.0, y + 0.28, 1.6, 0.35, model, font_size=15, color=c, bold=True)
    add_text_box(slide, 10.7, y + 0.08, 1.7, 0.5, val, font_size=14, color=LIGHT, alignment=PP_ALIGN.RIGHT)

# Future work
add_rounded_rect(slide, 0.8, 4.6, 5.6, 2.6, BG_CARD)
add_text_box(slide, 1.1, 4.7, 5, 0.4, "Future Work", font_size=18, color=ACCENT2, bold=True)
future = [
    "Evaluate on eval set with ground-truth annotations for proper F1/DER",
    "Alpha sweep (alpha = 0.3, 0.5, 0.7, 0.9) to balance KD vs hard loss",
    "Quantization (INT8/FP16) for additional 2-4x compression",
    "Latency benchmarking on CPU / ARM for edge deployment",
]
add_bullet_box(slide, 1.1, 5.2, 5.0, 2.2, future, font_size=14, color=LIGHT, bullet_color=ACCENT2)

# Experiment setup summary
add_rounded_rect(slide, 6.8, 4.6, 5.6, 2.6, BG_CARD)
add_text_box(slide, 7.1, 4.7, 5, 0.4, "Experiment Setup", font_size=18, color=ACCENT3, bold=True)
setup = [
    "Dataset: LibriParty (350 sessions, ~12 hrs audio)",
    "Teacher: SpeechBrain CRDNN (109K params)",
    "Training: 30 epochs, Adam, lr=1e-3, cosine schedule",
    "Hardware: NVIDIA RTX 3090, Windows 10",
]
add_bullet_box(slide, 7.1, 5.2, 5.0, 2.2, setup, font_size=14, color=LIGHT, bullet_color=ACCENT3)


# ── Save ──
output_path = "slides/VAD_KD_Presentation.pptx"
import os
os.makedirs("slides", exist_ok=True)
prs.save(output_path)
print(f"Saved to {output_path}")
